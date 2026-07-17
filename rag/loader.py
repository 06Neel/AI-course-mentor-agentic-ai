import os
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


@dataclass
class DocumentChunk:
    text: str
    source: str
    page: int
    chunk_id: int
    metadata: dict = field(default_factory=dict)


class DocumentLoader:
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def load_file(self, file_path: str) -> list[DocumentChunk]:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return self._load_pdf(file_path)
        elif ext == ".docx":
            return self._load_docx(file_path)
        elif ext == ".pptx":
            return self._load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _load_pdf(self, file_path: str) -> list[DocumentChunk]:
        reader = PdfReader(file_path)
        filename = os.path.basename(file_path)
        all_text = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                all_text.append((text.strip(), filename, i + 1))
        return self._create_chunks(all_text)

    def _load_docx(self, file_path: str) -> list[DocumentChunk]:
        doc = DocxDocument(file_path)
        filename = os.path.basename(file_path)
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if not full_text.strip():
            return []
        return self._create_chunks([(full_text, filename, 1)])

    def _load_pptx(self, file_path: str) -> list[DocumentChunk]:
        prs = Presentation(file_path)
        filename = os.path.basename(file_path)
        all_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                all_text.append(("\n".join(slide_text), filename, i + 1))
        return self._create_chunks(all_text)

    def _create_chunks(self, text_pages: list[tuple]) -> list[DocumentChunk]:
        chunks = []
        chunk_id = 0
        for text, source, page in text_pages:
            words = text.split()
            for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
                chunk_words = words[i : i + self.chunk_size]
                chunk_text = " ".join(chunk_words)
                if chunk_text.strip():
                    chunks.append(
                        DocumentChunk(
                            text=chunk_text,
                            source=source,
                            page=page,
                            chunk_id=chunk_id,
                            metadata={"source": source, "page": page},
                        )
                    )
                    chunk_id += 1
        return chunks

    def load_directory(self, dir_path: str) -> list[DocumentChunk]:
        all_chunks = []
        for f in Path(dir_path).iterdir():
            if f.suffix.lower() in [".pdf", ".docx", ".pptx"]:
                try:
                    chunks = self.load_file(str(f))
                    all_chunks.extend(chunks)
                except Exception as e:
                    print(f"Error loading {f.name}: {e}")
        return all_chunks
