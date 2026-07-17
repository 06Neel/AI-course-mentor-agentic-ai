import json
import uuid
from pathlib import Path
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from config.settings import Settings


@dataclass
class SlideContent:
    title: str
    bullets: list[str] = field(default_factory=list)
    notes: str = ""


@dataclass
class PPTContent:
    title: str
    subtitle: str = ""
    slides: list[SlideContent] = field(default_factory=list)


# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x0F, 0x34, 0x60)
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BULLET_COLOR = RGBColor(0xE0, 0xE0, 0xE0)


class PPTGenerator:
    def __init__(self):
        self.output_dir = Settings.DATA_DIR / "generated_ppts"
        self.output_dir.mkdir(exist_ok=True)

    def generate(self, content: PPTContent, filename: str = None) -> str:
        """Generate a PPT file and return its path."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, content.title, content.subtitle)

        for slide_data in content.slides:
            self._add_content_slide(prs, slide_data)

        self._add_closing_slide(prs, content.title)

        if not filename:
            safe_title = "".join(c if c.isalnum() else "_" for c in content.title)[:40]
            filename = f"{safe_title}_{uuid.uuid4().hex[:6]}.pptx"

        filepath = self.output_dir / filename
        prs.save(str(filepath))
        return str(filepath)

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._set_slide_bg(slide, DARK_BG)

        # Title
        left = Inches(1.5)
        top = Inches(2.2)
        width = Inches(10.333)
        height = Inches(1.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            top_sub = Inches(4.2)
            height_sub = Inches(1.0)
            txBox2 = slide.shapes.add_textbox(left, top_sub, width, height_sub)
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = SUBTITLE_GRAY
            p2.alignment = PP_ALIGN.CENTER

        # Accent line
        line_left = Inches(4.5)
        line_top = Inches(4.0)
        line_width = Inches(4.333)
        line = slide.shapes.add_shape(
            1, line_left, line_top, line_width, Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()

    def _add_content_slide(self, prs: Presentation, slide_data: SlideContent):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, DARK_BG)

        # Slide title
        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(11.733)
        height = Inches(1.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT

        # Accent line under title
        line = slide.shapes.add_shape(
            1, left, Inches(1.5), Inches(2.0), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()

        # Bullet points
        if slide_data.bullets:
            bullet_left = Inches(1.0)
            bullet_top = Inches(1.9)
            bullet_width = Inches(11.0)
            bullet_height = Inches(5.0)
            txBox2 = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)
            tf2 = txBox2.text_frame
            tf2.word_wrap = True

            for i, bullet in enumerate(slide_data.bullets):
                if i == 0:
                    p = tf2.paragraphs[0]
                else:
                    p = tf2.add_paragraph()
                p.text = f"  {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = BULLET_COLOR
                p.space_after = Pt(12)
                p.level = 0

        # Speaker notes
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.notes

    def _add_closing_slide(self, prs: Presentation, title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, DARK_BG)

        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(10.333)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        top2 = Inches(4.2)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(18)
        p2.font.color.rgb = SUBTITLE_GRAY
        p2.alignment = PP_ALIGN.CENTER

    def _set_slide_bg(self, slide, color: RGBColor):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def parse_llm_output(llm_text: str) -> PPTContent:
        """Parse LLM-generated JSON into PPTContent."""
        try:
            text = llm_text.strip()
            if text.startswith("```"):
                text = text.split("\n", 1)[1].rsplit("```", 1)[0].strip()
            data = json.loads(text)
            slides = []
            for s in data.get("slides", []):
                slides.append(SlideContent(
                    title=s.get("title", ""),
                    bullets=s.get("bullets", []),
                    notes=s.get("notes", ""),
                ))
            return PPTContent(
                title=data.get("title", "Presentation"),
                subtitle=data.get("subtitle", ""),
                slides=slides,
            )
        except (json.JSONDecodeError, KeyError, TypeError):
            return None

    @staticmethod
    def get_download_bytes(filepath: str) -> bytes:
        """Read a PPT file and return its bytes for download."""
        with open(filepath, "rb") as f:
            return f.read()
